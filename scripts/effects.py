#!/usr/bin/env python3
"""Post-processing of native OOXML effects that pptxgenjs cannot produce.

The pptxgenjs build tags the relevant shapes via their `objectName`; this script
reopens the .pptx with python-pptx, applies the matching native effect, then
rewrites the file. Without this pass, the shapes keep their solid fill (graceful
fallback) — so the deck stays valid even if the post-processing step is skipped.

Recognized markers (objectName):
  SOGRAD~<hex1>~<hex2>~<angle>          2-stop linear gradient (angle in degrees)
  SOGRAD~<hex1>~<hex2>~<hex3>~<angle>   3-stop linear gradient (0 / 50 / 100%)

Usage: python3 scripts/effects.py deck.pptx [deck_out.pptx]
"""
import sys
import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor


def _clean(hexs):
    return RGBColor.from_string(hexs.replace("#", "").strip()[:6])


def apply_gradient(shape, stops_hex, angle_deg):
    """Apply a native <a:gradFill> linear gradient while preserving the shadow/outline.
    Colors are validated BEFORE touching the fill: if a hex is invalid, we raise
    without having modified the shape (it keeps its solid fill -> graceful fallback upheld)."""
    colors = [_clean(h) for h in stops_hex]  # raises here if a hex is malformed, before any side effect
    fill = shape.fill
    fill.gradient()
    # angle: python-pptx counts counter-clockwise from left->right.
    try:
        fill.gradient_angle = float(angle_deg)
    except Exception:
        pass
    gs = fill.gradient_stops
    n = len(colors)
    # python-pptx creates 2 stops by default; adjust colors + positions.
    positions = [i / (n - 1) for i in range(n)] if n > 1 else [0.0, 1.0]
    for i, stop in enumerate(gs):
        if i < n:
            stop.color.rgb = colors[i]
            try:
                stop.position = positions[i]
            except Exception:
                pass
    # if the marker has more stops than the 2 created by default, inject the missing stops.
    if n > len(gs):
        gsLst = fill._xPr.find(qn("a:gradFill")).find(qn("a:gsLst"))
        while len(gsLst) < n:
            gsLst.append(copy.deepcopy(gsLst[-1]))
        for i, gs_el in enumerate(gsLst):
            srgb = gs_el.find(qn("a:srgbClr"))
            if srgb is not None and i < n:
                srgb.set("val", str(colors[i]))
            gs_el.set("pos", str(int(positions[i] * 100000)))


def process(path_in, path_out):
    prs = Presentation(path_in)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            name = (shape.name or "")
            if name.startswith("SOGRAD~"):
                parts = name.split("~")[1:]
                if len(parts) < 3:
                    continue
                angle = parts[-1]
                stops = parts[:-1]
                try:
                    apply_gradient(shape, stops, angle)
                    count += 1
                except Exception as e:
                    print(f"  ! gradient failed on '{name}': {e}", file=sys.stderr)
    prs.save(path_out)
    print(f"effects: {count} native gradient(s) applied -> {path_out}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: python3 scripts/effects.py deck.pptx [out.pptx]", file=sys.stderr)
        sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else inp
    process(inp, out)
