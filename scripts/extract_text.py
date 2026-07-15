#!/usr/bin/env python3
import signal, sys as _sys
try: signal.signal(signal.SIGPIPE, signal.SIG_DFL)  # clean exit when piped to head
except Exception: pass
"""Extract the text of a .pptx, one `## Slide N` section per slide (local equivalent of extract-text).

Usage: python3 scripts/extract_text.py deck.pptx
Includes table text (rows as `a | b | c`) and presenter notes (`[notes] ...`).
"""
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def shape_texts(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from shape_texts(sh.shapes)
        elif sh.has_text_frame and sh.text_frame.text.strip():
            yield sh.text_frame.text.strip()
        elif getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                yield " | ".join(c.text.strip() for c in row.cells)


prs = Presentation(sys.argv[1])
for i, slide in enumerate(prs.slides, 1):
    print(f"## Slide {i}")
    for t in shape_texts(slide.shapes):
        print(t)
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
        print(f"[notes] {slide.notes_slide.notes_text_frame.text.strip()}")
    print()
